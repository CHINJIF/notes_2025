from pdfminer.high_level import extract_text
from docx import Document
from pptx import Presentation
import argparse
import os

# what 
def output_to_txt(file_path, content):
    file_name = os.path.basename(file_path)
    output_path = os.path.join(os.path.dirname(file_path), f"{file_name}.txt")
    with open(output_path, 'w', encoding='utf-8') as txt_file:
        txt_file.write(content)

def pdf_to_text(file_path):
    # 提取 PDF 文件中的文本
    text = extract_text(file_path)
    
    output_to_txt(file_path, text)

def docx_to_text(file_path):
    # 提取 Word 文件中的文本
    doc = Document(file_path)
    text = '\n'.join([para.text for para in doc.paragraphs])
    
    output_to_txt(file_path, text)

def pptx_to_text(file_path):
    # 提取 PowerPoint 文件中的文本
    prs = Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    full_text = '\n'.join(text)
    
    output_to_txt(file_path, full_text)

def main():
    parser = argparse.ArgumentParser(description='Convert documents to text files.')
    parser.add_argument('doc_path', type=str, help='Path to the document file (PDF, DOCX, or PPTX).')
    parser.add_argument('doc_type', type=str, choices=['pdf', 'docx', 'pptx'], help='Type of the document (pdf, docx, or pptx).')
    args = parser.parse_args()

    if args.doc_type == 'pdf':
        pdf_to_text(args.doc_path)
    elif args.doc_type == 'docx':
        docx_to_text(args.doc_path)
    elif args.doc_type == 'pptx':
        pptx_to_text(args.doc_path)

if __name__ == "__main__":
    main()
